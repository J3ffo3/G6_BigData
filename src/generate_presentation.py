from __future__ import annotations

import json
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

PROJECT_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT_PATH = os.path.join(PROJECT_ROOT, "Presentacion_G6_BigData.pptx")
W5_FIGS = os.path.join(PROJECT_ROOT, "reports", "Week5", "figs")
W10_FIGS = os.path.join(PROJECT_ROOT, "reports", "Week10")

# ---- Palette ---------------------------------------------------------------
C_DARK   = RGBColor(0x1A, 0x23, 0x3E)   # dark navy
C_BLUE   = RGBColor(0x27, 0x5C, 0xAE)   # UPC-ish blue
C_ACCENT = RGBColor(0xE8, 0x6B, 0x2C)   # orange accent
C_LIGHT  = RGBColor(0xF2, 0xF5, 0xFB)   # slide background
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GRAY   = RGBColor(0x55, 0x5F, 0x6D)
C_GREEN  = RGBColor(0x2E, 0x7D, 0x32)

W, H = Cm(33.87), Cm(19.05)   # 16:9


# ---- Helpers ---------------------------------------------------------------

def blank_slide(prs: Presentation) -> object:
    layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(layout)


def bg(slide, color: RGBColor = C_LIGHT):
    from pptx.util import Emu
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h, text="", bold=False, size=18, color=C_DARK,
        bg_color=None, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = color
    if bg_color:
        from pptx.util import Pt as _Pt
        fill = txb.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
    return txb


def rect(slide, x, y, w, h, fill: RGBColor, line: RGBColor | None = None):
    from pptx.util import Pt as _Pt
    shp = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        x, y, w, h
    )
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp


def img(slide, path, x, y, w):
    if os.path.exists(path):
        slide.shapes.add_picture(path, x, y, width=w)


def header_bar(slide, title: str, subtitle: str = ""):
    rect(slide, 0, 0, W, Cm(2.6), C_DARK)
    box(slide, Cm(0.5), Cm(0.3), Cm(28), Cm(1.4),
        text=title, bold=True, size=22, color=C_WHITE)
    if subtitle:
        box(slide, Cm(0.5), Cm(1.6), Cm(28), Cm(0.9),
            text=subtitle, size=13, color=RGBColor(0xB0, 0xC4, 0xDE))
    rect(slide, 0, Cm(2.6), W, Cm(0.08), C_ACCENT)


def footer(slide, text="Group 6  |  Big Data  |  UPC"):
    rect(slide, 0, H - Cm(0.6), W, Cm(0.6), C_DARK)
    box(slide, Cm(0.5), H - Cm(0.58), W - Cm(1), Cm(0.55),
        text=text, size=8, color=RGBColor(0xAA, 0xBB, 0xCC))


def bullet_block(slide, items: list[str], x, y, w, h,
                 size=13, color=C_DARK, spacing=0.55):
    from pptx.util import Cm as _Cm
    for i, item in enumerate(items):
        box(slide, x, y + _Cm(i * spacing), w, _Cm(spacing + 0.1),
            text=f"•  {item}", size=size, color=color)


def colored_card(slide, x, y, w, h, header: str, body: str,
                 hdr_color=C_BLUE, body_size=11):
    rect(slide, x, y, w, Cm(0.7), hdr_color)
    box(slide, x + Cm(0.1), y, w - Cm(0.2), Cm(0.7),
        text=header, bold=True, size=11, color=C_WHITE)
    rect(slide, x, y + Cm(0.7), w, h - Cm(0.7), C_WHITE,
         line=RGBColor(0xCC, 0xCC, 0xCC))
    box(slide, x + Cm(0.15), y + Cm(0.75), w - Cm(0.3), h - Cm(0.9),
        text=body, size=body_size, color=C_GRAY)


# ============================================================================
# SLIDES
# ============================================================================

def slide_title(prs):
    s = blank_slide(prs)
    bg(s, C_DARK)
    rect(s, 0, Cm(5.5), W, Cm(0.18), C_ACCENT)

    box(s, Cm(1.5), Cm(1.2), Cm(30), Cm(2),
        text="Semester Project: Domain Discovery,", bold=False, size=17,
        color=RGBColor(0xB0, 0xC4, 0xDE), align=PP_ALIGN.CENTER)
    box(s, Cm(1.5), Cm(2.7), Cm(30), Cm(2),
        text="Recommendation & Graph Intelligence", bold=False, size=17,
        color=RGBColor(0xB0, 0xC4, 0xDE), align=PP_ALIGN.CENTER)

    box(s, Cm(1.5), Cm(5.8), Cm(30), Cm(2.5),
        text="Week 10: Recommendation Engine", bold=True, size=30,
        color=C_WHITE, align=PP_ALIGN.CENTER)

    box(s, Cm(1.5), Cm(9.2), Cm(30), Cm(1.2),
        text="Group 6  —  Big Data Course  —  UPC", size=15,
        color=C_ACCENT, align=PP_ALIGN.CENTER)

    members = ("Iam Alvarez Orellana   •   Jeffrey Diaz Villanueva\n"
               "Paula Mancilla Cienfuegos   •   Fernando Paredes Espinoza")
    box(s, Cm(2), Cm(11.2), Cm(30), Cm(2),
        text=members, size=12, color=RGBColor(0x8A, 0xA8, 0xCC),
        align=PP_ALIGN.CENTER)


def slide_pipeline(prs):
    s = blank_slide(prs)
    bg(s, C_LIGHT)
    header_bar(s, "Project Pipeline", "4-layer architecture from raw data to recommendations")
    footer(s)

    layers = [
        ("W3  DATA", "steam_v1.parquet\n480 k interactions\nStar schema on AppID", C_BLUE),
        ("W5  FEATURES", "X_numeric (87k×26)\nX_categorical (87k×5k)\nX_text TF-IDF (87k×5k)", C_BLUE),
        ("W7  CLUSTERS", "Tags-only binary matrix\nK-Means K=6 (SVD+L2)\n6 semantic segments", C_BLUE),
        ("W10  RECS", "Popularity baseline\nContent + MF + Hybrid\nHR@10 evaluation", C_ACCENT),
    ]

    xs = [Cm(0.6), Cm(9.1), Cm(17.6), Cm(26.1)]
    for i, (lbl, body, col) in enumerate(layers):
        colored_card(s, xs[i], Cm(3.2), Cm(7.2), Cm(10.5), lbl, body, col)
        if i < 3:
            box(s, xs[i] + Cm(7.3), Cm(7.5), Cm(1.6), Cm(1),
                text="→", bold=True, size=28, color=C_BLUE, align=PP_ALIGN.CENTER)


def slide_dataset(prs):
    s = blank_slide(prs)
    bg(s, C_LIGHT)
    header_bar(s, "Dataset", "Steam Games & Reviews — Kaggle CC0")
    footer(s)

    stats = [
        ("80,000+", "Steam games in full catalog"),
        ("480,025", "Interaction records (V1 subset, 500 games)"),
        ("426,601", "Unique users"),
        ("87,806", "Games in feature matrices"),
        ("69,943", "Tagged games (used for clustering)"),
        ("99.8 %", "Interaction matrix sparsity"),
    ]
    for i, (val, desc) in enumerate(stats):
        col = 0 if i < 3 else 1
        row = i % 3
        x = Cm(1.0) if col == 0 else Cm(17.5)
        y = Cm(3.2) + Cm(row * 4.5)
        rect(s, x, y, Cm(14.5), Cm(3.8), C_WHITE, line=RGBColor(0xCC, 0xCC, 0xCC))
        rect(s, x, y, Cm(14.5), Cm(1.6), C_BLUE)
        box(s, x + Cm(0.2), y + Cm(0.1), Cm(14), Cm(1.4),
            text=val, bold=True, size=26, color=C_WHITE)
        box(s, x + Cm(0.2), y + Cm(1.7), Cm(14), Cm(1.9),
            text=desc, size=12, color=C_GRAY)

    box(s, Cm(1), Cm(17.0), Cm(32), Cm(0.8),
        text="Star schema: AppID joins Games catalog ↔ Reviews interactions",
        size=11, color=C_GRAY, align=PP_ALIGN.CENTER)


def slide_clustering_theory(prs):
    s = blank_slide(prs)
    bg(s, C_LIGHT)
    header_bar(s, "Clustering — Class 7 Methodology",
               "How lecture concepts drove implementation decisions")
    footer(s)

    decisions = [
        ("Lloyd’s Algorithm\nVerification",
         "Class 7 defines K-means as\niterative assignment + centroid\nupdate. We manually ran one\niteration per seed to confirm\nJ_after ≤ J_before (5/5 seeds).",
         C_BLUE),
        ("Representation\nAblation",
         "Class 7: distance metric depends\non representation. We tested\n3 variants:\nRaw sparse → sil=0.026\nSVD 50D → sil=0.064\nSVD+L2 → sil=0.077 ✓",
         C_BLUE),
        ("K Selection\n(Elbow + Silhouette)",
         "Class 7: no universal K.\nWe swept K∈{4,6,8,10}.\nK=6 chosen: balanced\ncluster sizes (min 4,773)\nvs K=8/10 (min 958 games).",
         C_ACCENT),
    ]

    xs = [Cm(0.6), Cm(11.6), Cm(22.6)]
    for i, (hdr, body, col) in enumerate(decisions):
        colored_card(s, xs[i], Cm(3.0), Cm(10.5), Cm(13.5), hdr, body, col, body_size=12)


def slide_cluster_results(prs):
    s = blank_slide(prs)
    bg(s, C_LIGHT)
    header_bar(s, "K=6 Cluster Profiles", "Top tags by lift = P(tag|cluster) / P(tag)")
    footer(s)

    clusters = [
        ("C0  Creative Tools",   "4,773  (6.8%)",  "video production 14.6×\naudio production 14.6×\nphoto editing 14.6×",       C_GRAY),
        ("C1  Action",           "15,353 (22%)",   "combat 3.2×\naction-adventure 3.1×\nside scroller 3.1×",                      C_BLUE),
        ("C2  Narrative",        "13,576 (19.4%)", "text-based 4.0×\nstory rich 3.1×\ngreat soundtrack 2.9×",                      C_BLUE),
        ("C3  Strategy/Puzzle",  "14,060 (20.1%)", "card game 3.8×\nlogic 3.7×\nbuilding 3.3×",                                   C_BLUE),
        ("C4  Local Multiplayer","12,112 (17.3%)", "local co-op 5.8×\n4 player local 5.8×\nsplit screen 5.7×",                    C_ACCENT),
        ("C5  Casual/HO",        "10,069 (14.4%)", "hidden object 6.5×\npoint & click 5.5×",                                           C_BLUE),
    ]

    xs = [Cm(0.5), Cm(11.5), Cm(22.5)]
    ys = [Cm(3.0), Cm(10.6)]
    for i, (name, size, tags, col) in enumerate(clusters):
        x = xs[i % 3]
        y = ys[i // 3]
        colored_card(s, x, y, Cm(10.2), Cm(6.8), name, f"Size: {size}\n\n{tags}", col, body_size=11)

    box(s, Cm(0.5), H - Cm(1.5), Cm(33), Cm(0.8),
        text="Silhouette = 0.077  |  Inertia = 29,175  |  ARI stability mean = 0.49",
        size=11, color=C_GRAY, align=PP_ALIGN.CENTER)


def slide_rec_architecture(prs):
    s = blank_slide(prs)
    bg(s, C_LIGHT)
    header_bar(s, "Recommendation Engine — Architecture (Week 10)",
               "4 systems evaluated on the same 484-game universe")
    footer(s)

    systems = [
        ("1. Popularity\nBaseline",
         "Rank by aggregate\n'recommendations' count.\nSame list for every user.\nNo personalization.",
         C_GRAY),
        ("2. Content-Based\nBaseline",
         "TF-IDF on name+tags\n+description (3k terms).\nSVD → 50D. User vector\n= centroid of train games.\nCosine similarity ranking.",
         C_BLUE),
        ("3. Matrix\nFactorization (SGD)",
         "Binary interaction matrix\nR ∈ {0,1}^{users × games}.\nLatent factors P, Q.\nObjective: min MSE + L2 reg.\nk ∈ {10, 20, 50} swept.",
         C_BLUE),
        ("4. Hybrid\nRecommender",
         "α × content_score\n+ (1-α) × MF_score\nBoth min-max normalized.\nα ∈ {0.1,0.3,0.5,0.7,0.9}\nBest: α = 0.3",
         C_ACCENT),
    ]

    xs = [Cm(0.5), Cm(9.0), Cm(17.5), Cm(26.0)]
    for i, (hdr, body, col) in enumerate(systems):
        colored_card(s, xs[i], Cm(3.0), Cm(7.8), Cm(12.5), hdr, body, col, body_size=11)

    box(s, Cm(0.5), Cm(16.0), Cm(33), Cm(1.0),
        text="AppID is the common key across catalog, interaction parquet, and cluster labels — hybrid combination is exact by construction.",
        size=11, color=C_GRAY, align=PP_ALIGN.CENTER)


def slide_mf_detail(prs):
    s = blank_slide(prs)
    bg(s, C_LIGHT)
    header_bar(s, "Matrix Factorization — SGD", "Objective function and update rules")
    footer(s)

    rect(s, Cm(0.6), Cm(3.0), Cm(32.7), Cm(7.5), C_WHITE,
         line=RGBColor(0xCC, 0xCC, 0xCC))
    box(s, Cm(1.0), Cm(3.2), Cm(32), Cm(1.0),
        text="Objective:", bold=True, size=14, color=C_DARK)
    box(s, Cm(1.5), Cm(4.0), Cm(32), Cm(1.4),
        text="min_{P,Q}  Σ_{(u,i)∈Ω}  (R_ui - p_uᵀ q_i)²  +  λ ( ‖P‖² + ‖Q‖² )",
        size=16, color=C_BLUE)
    box(s, Cm(1.0), Cm(5.5), Cm(32), Cm(1.0),
        text="SGD updates per observed pair (u, i):", bold=True, size=14, color=C_DARK)
    box(s, Cm(1.5), Cm(6.3), Cm(32), Cm(0.9),
        text="p_u  ←  p_u + η ( e_ui · q_i  −  λ · p_u )     where  e_ui = R_ui − p_uᵀ q_i",
        size=14, color=C_BLUE)
    box(s, Cm(1.5), Cm(7.1), Cm(32), Cm(0.9),
        text="q_i  ←  q_i + η ( e_ui · p_u  −  λ · q_i )",
        size=14, color=C_BLUE)

    params = [
        ("Learning rate η", "0.03", "Standard for implicit binary data"),
        ("Regularization λ", "0.01", "Prevents overfitting (99.8% sparse matrix)"),
        ("Epochs", "30", "Convergence confirmed by loss curve"),
        ("Latent dim k", "{10, 20, 50}", "Swept; best k=10 on 484-game universe"),
        ("Initialization", "N(0, 0.01)", "Small values break symmetry"),
    ]
    box(s, Cm(0.6), Cm(10.8), Cm(32), Cm(0.8),
        text="Hyperparameters:", bold=True, size=13, color=C_DARK)
    for i, (param, val, note) in enumerate(params):
        y = Cm(11.6) + Cm(i * 1.2)
        rect(s, Cm(0.6), y, Cm(32.7), Cm(1.1), C_LIGHT if i % 2 == 0 else C_WHITE,
             line=RGBColor(0xDD, 0xDD, 0xDD))
        box(s, Cm(0.8), y + Cm(0.1), Cm(9), Cm(0.9), text=param, bold=True, size=11, color=C_DARK)
        box(s, Cm(9.8), y + Cm(0.1), Cm(5), Cm(0.9), text=val, size=11, color=C_ACCENT)
        box(s, Cm(14.8), y + Cm(0.1), Cm(18), Cm(0.9), text=note, size=11, color=C_GRAY)


def slide_eval_protocol(prs):
    s = blank_slide(prs)
    bg(s, C_LIGHT)
    header_bar(s, "Evaluation Protocol", "Hit-Rate@K on held-out positive interactions")
    footer(s)

    steps = [
        ("1", "Build interaction matrix", "Sparse binary R (users × games). voted_up=True → 1. Duplicates removed."),
        ("2", "Qualify users", "Keep users with ≥ 2 positive interactions (need ≥1 train + 1 hold-out)."),
        ("3", "Train/test split", "Hold out 1 random positive per user (random_state=42). Rest = training."),
        ("4", "Candidate pool", "All games the user has NOT interacted with in training."),
        ("5", "Compute HR@K", "Fraction of users whose held-out game appears in top-K.\nPrimary metric: HR@10."),
        ("6", "Leakage prevention", "Training items masked to −∞ before ranking. Hold-out never used in training."),
    ]

    for i, (num, title, body) in enumerate(steps):
        col = 0 if i < 3 else 1
        row = i % 3
        x = Cm(0.5) if col == 0 else Cm(17.2)
        y = Cm(3.0) + Cm(row * 4.8)
        rect(s, x, y, Cm(3.0), Cm(4.2), C_BLUE)
        box(s, x, y + Cm(0.8), Cm(3.0), Cm(2.0), text=num, bold=True, size=36,
            color=C_WHITE, align=PP_ALIGN.CENTER)
        rect(s, x + Cm(3.0), y, Cm(13.0), Cm(4.2), C_WHITE,
             line=RGBColor(0xCC, 0xCC, 0xCC))
        box(s, x + Cm(3.1), y + Cm(0.2), Cm(12.7), Cm(1.0),
            text=title, bold=True, size=12, color=C_DARK)
        box(s, x + Cm(3.1), y + Cm(1.2), Cm(12.7), Cm(2.8),
            text=body, size=11, color=C_GRAY)

    box(s, Cm(0.5), H - Cm(1.5), Cm(33), Cm(0.8),
        text="12,389 qualifying users  |  484-game universe  |  426,601 total users",
        size=11, color=C_GRAY, align=PP_ALIGN.CENTER)


def slide_results_chart(prs):
    s = blank_slide(prs)
    bg(s, C_LIGHT)
    header_bar(s, "Results — HR@K Comparison", "All 4 systems evaluated at K ∈ {5, 10, 20}")
    footer(s)

    fig = os.path.join(W10_FIGS, "fig01_hr_comparison.png")
    img(s, fig, Cm(0.8), Cm(3.0), Cm(21))

    results = [
        ("Popularity",         "0.4597", "0.5796", "0.7124", C_GRAY),
        ("Content-based",      "0.1806", "0.2170", "0.2669", C_BLUE),
        ("MF k=10",            "0.1103", "0.1888", "0.2949", C_BLUE),
        ("MF k=20",            "0.0764", "0.1333", "0.2183", C_BLUE),
        ("MF k=50",            "0.0756", "0.1235", "0.1998", C_BLUE),
        ("Hybrid α=0.3 ★", "0.2057", "0.2569", "0.3330", C_ACCENT),
    ]

    x0 = Cm(22.5)
    box(s, x0, Cm(3.0), Cm(11), Cm(0.8), text="System", bold=True, size=10, color=C_DARK)
    box(s, x0 + Cm(7.0), Cm(3.0), Cm(2), Cm(0.8), text="HR@5", bold=True, size=10, color=C_DARK)
    box(s, x0 + Cm(8.7), Cm(3.0), Cm(2), Cm(0.8), text="HR@10", bold=True, size=10, color=C_DARK)
    box(s, x0 + Cm(10.4), Cm(3.0), Cm(2), Cm(0.8), text="HR@20", bold=True, size=10, color=C_DARK)

    for i, (name, h5, h10, h20, col) in enumerate(results):
        y = Cm(3.9) + Cm(i * 1.9)
        fill = C_LIGHT if i % 2 == 0 else C_WHITE
        rect(s, x0, y, Cm(11.2), Cm(1.8), fill, line=RGBColor(0xDD, 0xDD, 0xDD))
        box(s, x0 + Cm(0.1), y + Cm(0.3), Cm(6.8), Cm(1.2), text=name,
            bold=(i == 5), size=11, color=col)
        for j, val in enumerate([h5, h10, h20]):
            box(s, x0 + Cm(7.0 + j * 1.7), y + Cm(0.3), Cm(1.6), Cm(1.2),
                text=val, bold=(i == 5), size=11,
                color=C_ACCENT if i == 5 else C_DARK, align=PP_ALIGN.CENTER)


def slide_mf_charts(prs):
    s = blank_slide(prs)
    bg(s, C_LIGHT)
    header_bar(s, "MF Training & K-Sweep Analysis", "Convergence and latent dimension selection")
    footer(s)

    img(s, os.path.join(W10_FIGS, "fig02_mf_loss_curve.png"), Cm(0.5), Cm(3.0), Cm(16.5))
    img(s, os.path.join(W10_FIGS, "fig03_mf_k_sweep.png"),   Cm(17.2), Cm(3.0), Cm(16.0))

    box(s, Cm(0.5), Cm(14.5), Cm(16.5), Cm(3.5),
        text=("Loss decreases monotonically over 30 epochs (k=10). "
              "Convergence confirmed — regularized MSE drops from 0.98 to 0.016."),
        size=12, color=C_GRAY)
    box(s, Cm(17.2), Cm(14.5), Cm(16.0), Cm(3.5),
        text=("k=10 beats k=20 and k=50 on HR@10. "
              "Smaller latent space generalizes better on the sparse 484-game universe. "
              "Higher k overfits the scarce interaction signal."),
        size=12, color=C_GRAY)


def slide_hybrid_chart(prs):
    s = blank_slide(prs)
    bg(s, C_LIGHT)
    header_bar(s, "Hybrid System — Alpha Sweep", "Balancing content signal vs collaborative filtering")
    footer(s)

    img(s, os.path.join(W10_FIGS, "fig04_hybrid_alpha_sweep.png"), Cm(0.5), Cm(3.0), Cm(18))

    box(s, Cm(19.5), Cm(3.5), Cm(14), Cm(2.2),
        text="s_hybrid(u,g) =", bold=True, size=15, color=C_DARK)
    box(s, Cm(19.5), Cm(5.2), Cm(14), Cm(2.2),
        text="α × s_content(u,g)", size=14, color=C_BLUE)
    box(s, Cm(19.5), Cm(6.6), Cm(14), Cm(1.5),
        text="+ (1-α) × s_MF(u,g)", size=14, color=C_GREEN)

    findings = [
        "α = 0.3 is optimal → HR@10 = 0.2569",
        "30% content + 70% MF",
        "Beats content-only (0.217) by +18%",
        "Beats MF-only (0.189) by +36%",
        "Both signals needed — neither dominates",
    ]
    bullet_block(s, findings, Cm(19.5), Cm(9.5), Cm(14), Cm(7), size=13)


def slide_pca_figure(prs):
    s = blank_slide(prs)
    bg(s, C_LIGHT)
    header_bar(s, "Feature Engineering (W5) — Dimensionality Reduction",
               "PCA on X_numeric and TruncatedSVD on X_text")
    footer(s)

    img(s, os.path.join(W5_FIGS, "fig_w5_01_pca_variance.png"), Cm(0.5), Cm(3.0), Cm(16))
    img(s, os.path.join(W5_FIGS, "fig_w5_02_svd_spectrum.png"), Cm(17.0), Cm(3.0), Cm(16.5))

    box(s, Cm(0.5), Cm(14.5), Cm(16), Cm(2.5),
        text="PCA: k=12 components capture ≥90% of numeric variance.\n"
             "12D used for Lloyd’s verification in W7.",
        size=12, color=C_GRAY)
    box(s, Cm(17.0), Cm(14.5), Cm(16.5), Cm(2.5),
        text="SVD on X_text: slow singular value decay (broad vocabulary).\n"
             "k=201 for 50% variance. k=50 used in clustering (tag matrix).",
        size=12, color=C_GRAY)


def slide_conclusions(prs):
    s = blank_slide(prs)
    bg(s, C_LIGHT)
    header_bar(s, "Key Findings & Conclusions")
    footer(s)

    findings = [
        ("Class 7 methodology validated",
         "Lloyd’s descent confirmed on 5/5 seeds. SVD+L2 representation "
         "improved silhouette from 0.026 (raw) to 0.077. K=6 chosen via "
         "balanced elbow + cluster size analysis.",
         C_BLUE),
        ("Recommendation engine delivers",
         "Hybrid α=0.3 achieves HR@10=0.2569 — best among personalized systems. "
         "Content-based alone (0.217) and MF alone (0.189) are each sub-optimal. "
         "Combining both signals matters.",
         C_ACCENT),
        ("Popularity ceiling is high",
         "Popularity HR@10=0.58 is dominant because the evaluation universe is "
         "small (484 games). Full-catalog evaluation would likely reverse this "
         "result as user tastes diverge from average.",
         C_GRAY),
    ]

    for i, (title, body, col) in enumerate(findings):
        y = Cm(3.2) + Cm(i * 4.8)
        rect(s, Cm(0.5), y, Cm(0.4), Cm(4.2), col)
        box(s, Cm(1.2), y + Cm(0.2), Cm(32), Cm(1.0),
            text=title, bold=True, size=14, color=col)
        box(s, Cm(1.2), y + Cm(1.2), Cm(32), Cm(2.8),
            text=body, size=12, color=C_GRAY)


def slide_pipeline_commands(prs):
    s = blank_slide(prs)
    bg(s, C_DARK)
    header_bar(s, "Reproducible Pipeline", "All results generated from 5 commands")
    footer(s)

    cmds = [
        ("Step 1  —  Ingestion",          "python src/ingest.py",              "steam_v1.parquet → data/processed/"),
        ("Step 2  —  Feature Engineering", "python src/feature_engineering.py", "X_numeric, X_categorical, X_text → artifacts/"),
        ("Step 3  —  Clustering",          "python src/clustering_week7.py",    "cluster_labels_k6.npy → artifacts/clustering/"),
        ("Step 4  —  Recommendation",      "python src/recommendation.py",      "rec_metrics.json → artifacts/recommendation/"),
        ("Step 5  —  Figures & Reports",   "python src/generate_figures.py",    "PNG figures → reports/Week5/ and reports/Week10/"),
    ]

    for i, (label, cmd, output) in enumerate(cmds):
        y = Cm(3.0) + Cm(i * 3.0)
        rect(s, Cm(0.5), y, Cm(8), Cm(2.6), C_BLUE)
        box(s, Cm(0.6), y + Cm(0.6), Cm(7.8), Cm(1.4),
            text=label, bold=True, size=11, color=C_WHITE)
        rect(s, Cm(8.7), y, Cm(14), Cm(2.6), RGBColor(0x0D, 0x14, 0x2B))
        box(s, Cm(8.8), y + Cm(0.6), Cm(13.8), Cm(1.4),
            text=cmd, size=13, color=RGBColor(0x7D, 0xFF, 0xA0))
        rect(s, Cm(23.0), y, Cm(10.5), Cm(2.6), RGBColor(0x12, 0x1A, 0x35))
        box(s, Cm(23.1), y + Cm(0.6), Cm(10.2), Cm(1.4),
            text=output, size=11, color=RGBColor(0xAA, 0xBB, 0xCC))

    box(s, Cm(0.5), H - Cm(1.5), Cm(33), Cm(0.8),
        text="random_state=42 throughout  —  all results fully deterministic",
        size=11, color=RGBColor(0x88, 0xAA, 0xCC), align=PP_ALIGN.CENTER)


# ============================================================================
# MAIN
# ============================================================================

def build():
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    slide_title(prs)
    slide_pipeline(prs)
    slide_dataset(prs)
    slide_pca_figure(prs)
    slide_clustering_theory(prs)
    slide_cluster_results(prs)
    slide_rec_architecture(prs)
    slide_mf_detail(prs)
    slide_eval_protocol(prs)
    slide_results_chart(prs)
    slide_mf_charts(prs)
    slide_hybrid_chart(prs)
    slide_conclusions(prs)
    slide_pipeline_commands(prs)

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
